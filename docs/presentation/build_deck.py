#!/usr/bin/env python3
"""
build_deck.py — generate docs/presentation/valkyrie-overview.pptx.

A visual-first technical-team deck. Run:
    pip3 install --user python-pptx
    python3 docs/presentation/build_deck.py
Output lands at docs/presentation/valkyrie-overview.pptx.

The deck is regenerable — edit this file, not the .pptx, so the source of
truth lives in git.
"""

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ─── palette ──────────────────────────────────────────────────────────────
NAVY      = RGBColor(0x0a, 0x0a, 0x12)
NAVY_DEEP = RGBColor(0x05, 0x05, 0x10)
PANEL     = RGBColor(0x16, 0x16, 0x22)
CYAN      = RGBColor(0x22, 0xd3, 0xee)
CYAN_DIM  = RGBColor(0x0e, 0x74, 0x90)
AMBER     = RGBColor(0xf5, 0x9e, 0x0b)
AMBER_DIM = RGBColor(0x78, 0x4e, 0x05)
EMERALD   = RGBColor(0x10, 0xb9, 0x81)
VIOLET    = RGBColor(0x8b, 0x5c, 0xf6)
ROSE      = RGBColor(0xf4, 0x3f, 0x5e)
WHITE     = RGBColor(0xff, 0xff, 0xff)
WHITE_DIM = RGBColor(0xb8, 0xb8, 0xc8)
WHITE_FAINT = RGBColor(0x6a, 0x6a, 0x80)
GRAY      = RGBColor(0x2a, 0x2a, 0x38)
BORDER    = RGBColor(0x33, 0x33, 0x42)

# ─── deck setup ───────────────────────────────────────────────────────────
prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


def add_slide():
    s = prs.slides.add_slide(BLANK)
    bg = s.background.fill
    bg.solid()
    bg.fore_color.rgb = NAVY
    return s


def add_text(s, x, y, w, h, text, *, size=18, color=WHITE, bold=False,
             align=PP_ALIGN.LEFT, italic=False, font="Helvetica Neue",
             anchor=MSO_ANCHOR.TOP):
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Inches(0)
    tf.margin_top = tf.margin_bottom = Inches(0)
    lines = text.split('\n')
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run()
        r.text = line
        r.font.size = Pt(size)
        r.font.color.rgb = color
        r.font.bold = bold
        r.font.italic = italic
        r.font.name = font
    return tb


def add_bullets(s, x, y, w, h, items, *, size=16, color=WHITE_DIM,
                bullet=" ▸ ", line_spacing=1.25):
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0)
    tf.margin_top = tf.margin_bottom = Inches(0)
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = line_spacing
        # bullet symbol in accent color
        rb = p.add_run()
        rb.text = bullet
        rb.font.size = Pt(size)
        rb.font.color.rgb = CYAN
        rb.font.bold = True
        rb.font.name = "Helvetica Neue"
        # body
        r = p.add_run()
        r.text = item
        r.font.size = Pt(size)
        r.font.color.rgb = color
        r.font.name = "Helvetica Neue"
    return tb


def add_block(s, x, y, w, h, label, *, fill=CYAN, text=NAVY, size=16,
              bold=True, shape=MSO_SHAPE.ROUNDED_RECTANGLE, line=None,
              line_w=1.0, font="Helvetica Neue"):
    sh = s.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    sh.line.color.rgb = line if line is not None else fill
    sh.line.width = Pt(line_w)
    # adjust rounded corner radius — smaller radius for blocks
    try:
        sh.adjustments[0] = 0.15
    except Exception:
        pass
    tf = sh.text_frame
    tf.margin_left = Inches(0.08)
    tf.margin_right = Inches(0.08)
    tf.margin_top = Inches(0.06)
    tf.margin_bottom = Inches(0.06)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = label
    r.font.size = Pt(size)
    r.font.color.rgb = text
    r.font.bold = bold
    r.font.name = font
    return sh


def add_outline(s, x, y, w, h, label, *, color=CYAN, size=14,
                line_w=1.5, dash=None, shape=MSO_SHAPE.ROUNDED_RECTANGLE):
    """Hollow / outline block."""
    sh = s.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    sh.fill.background()  # transparent fill
    sh.line.color.rgb = color
    sh.line.width = Pt(line_w)
    if dash is not None:
        sh.line.dash_style = dash
    try:
        sh.adjustments[0] = 0.08
    except Exception:
        pass
    if label:
        tf = sh.text_frame
        tf.margin_left = Inches(0.08)
        tf.margin_right = Inches(0.08)
        tf.margin_top = Inches(0.04)
        tf.vertical_anchor = MSO_ANCHOR.TOP
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        r = p.add_run()
        r.text = label
        r.font.size = Pt(size)
        r.font.color.rgb = color
        r.font.bold = True
        r.font.name = "Helvetica Neue"
    return sh


def add_arrow(s, x1, y1, x2, y2, *, color=CYAN, line_w=2.5):
    """Straight arrow connector, single-headed."""
    conn = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                                  Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    conn.line.color.rgb = color
    conn.line.width = Pt(line_w)
    # arrowhead at end (use lxml to set tailEnd/headEnd)
    from pptx.oxml.ns import qn
    ln = conn.line._get_or_add_ln()
    # remove existing endings
    for tag in ('a:headEnd', 'a:tailEnd'):
        for el in ln.findall(qn(tag)):
            ln.remove(el)
    from lxml import etree
    nsmap = {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'}
    tail = etree.SubElement(ln, '{http://schemas.openxmlformats.org/drawingml/2006/main}tailEnd')
    tail.set('type', 'triangle')
    tail.set('w', 'med')
    tail.set('len', 'med')
    return conn


def add_title_bar(s, title, kicker=None):
    """Header bar at top of slide."""
    # Thin cyan rule
    rule = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(0.55), Inches(0.5), Inches(0.04))
    rule.fill.solid(); rule.fill.fore_color.rgb = CYAN
    rule.line.fill.background()
    if kicker:
        add_text(s, 1.25, 0.42, 8, 0.3, kicker.upper(), size=10, bold=True,
                 color=CYAN, font="Helvetica Neue")
    add_text(s, 0.6, 0.68, 12, 0.7, title, size=30, bold=True, color=WHITE,
             font="Helvetica Neue")


def add_footer(s, txt, idx, total):
    add_text(s, 0.6, 7.05, 8, 0.3, txt, size=10, color=WHITE_FAINT, italic=True)
    add_text(s, 12.0, 7.05, 1.0, 0.3, f"{idx} / {total}", size=10,
             color=WHITE_FAINT, align=PP_ALIGN.RIGHT)


# ─── slides ───────────────────────────────────────────────────────────────

slide_count = 28  # update as needed; used for footer

# 1. TITLE ─────────────────────────────────────────────────────────────────
s = add_slide()
# big rule
r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(2.7), Inches(0.6), Inches(0.08))
r.fill.solid(); r.fill.fore_color.rgb = CYAN; r.line.fill.background()
add_text(s, 0.8, 2.0, 8, 0.5, "VALKYRIE", size=14, bold=True, color=CYAN)
add_text(s, 0.8, 2.85, 12, 1.3, "Disciplined Agentic Engineering", size=54,
         bold=True, color=WHITE)
add_text(s, 0.8, 4.05, 12, 0.8, "on Claude Code — built on Matt Pocock's skills.",
         size=28, color=WHITE_DIM)
add_text(s, 0.8, 5.4, 12, 0.6,
         "Why the workflow is the product, not the model.", size=18,
         color=WHITE_FAINT, italic=True)
add_text(s, 0.8, 6.7, 12, 0.3,
         "github.com/BunnyDAO/Valkyrie", size=11, color=CYAN_DIM)

# 2. THE PROBLEM ─────────────────────────────────────────────────────────
s = add_slide()
add_title_bar(s, "Raw Claude Code is fast and structureless.", kicker="The problem")
add_bullets(s, 0.7, 1.7, 12.0, 5.0, [
    "The model decides when to write code — usually too early.",
    "Long sessions drift: context shrinks, the early decisions get forgotten.",
    "Context-limit handoff resets the orchestrator and the human-in-the-loop gates with it.",
    "Each answer is a wall of prose; you have to mine the next decision out of it by hand.",
    "Nothing enforces \"design before code,\" \"PRD before issues,\" or \"approval before commit.\"",
    "Cost ladders up the wrong way: opus burned on chat, sonnet asked to write tests it doesn't understand.",
], size=20, line_spacing=1.4)
add_footer(s, "The problem Valkyrie was built to fix.", 2, slide_count)

# 3. MATT POCOCK INTRO ────────────────────────────────────────────────────
s = add_slide()
add_title_bar(s, "Matt Pocock's `skills` repo — the foundation.", kicker="The base layer")
add_text(s, 0.7, 1.65, 12.0, 0.5,
         "github.com/mattpocock/skills — open-source Claude Code skill prompts.",
         size=18, color=CYAN, italic=True)
add_bullets(s, 0.7, 2.3, 12.0, 4.2, [
    "Structured `SKILL.md` prompts that route a request through stages instead of jumping to code.",
    "Each skill carries its own discipline: a grilling pattern, a PRD template, an issue-slicing rule.",
    "Plus the \"Ralph\" pattern (aihero.dev): autonomous loops that re-spawn fresh CLI sessions per issue.",
    "Pocock = the *content* of the workflow.  Valkyrie = the *enforcement, ergonomics, and observability around it*.",
], size=18, line_spacing=1.45)
add_text(s, 0.7, 6.5, 12.0, 0.4,
         "Six of our ten skills are direct adaptations of mattpocock/skills — credited in every SKILL.md header.",
         size=12, color=WHITE_FAINT, italic=True)
add_footer(s, "The base layer we build on.", 3, slide_count)

# 4. DESIGN (grill-with-docs) ────────────────────────────────────────────
def stage_slide(idx, kicker, title, leftblock, leftblock_col, rightbullets, footer):
    s = add_slide()
    add_title_bar(s, title, kicker=kicker)
    # left big stage block
    add_block(s, 0.7, 1.9, 4.4, 4.0, leftblock, fill=leftblock_col,
              text=NAVY, size=28, bold=True)
    # right bullets
    add_bullets(s, 5.6, 2.0, 7.4, 4.5, rightbullets, size=18,
                line_spacing=1.45)
    add_footer(s, footer, idx, slide_count)
    return s

stage_slide(
    4, "Stage 1 — Pocock",
    "DESIGN  —  the grilling.",
    "DESIGN\n(grill-with-docs)",
    AMBER,
    [
        "Interview-driven: one question at a time, no monologue.",
        "Grills against scope, actors, failure modes, data model, boundaries, observability.",
        "Codebase exploration delegated to single-task sub-agents — keeps the main thread lean.",
        "Updates the glossary inline (`CONTEXT.md`) and writes ADRs only when the decision is load-bearing.",
        "No production code at this stage — the TDD gate enforces it mechanically.",
    ],
    "DESIGN — front-loaded thinking, before any code can be written.",
)

# 5. PRD (to-prd) ─────────────────────────────────────────────────────────
stage_slide(
    5, "Stage 2 — Pocock",
    "PRD  —  the contract.",
    "PRD\n(to-prd)",
    AMBER,
    [
        "Synthesizes the grilling into one decision-ready document.",
        "Implementation decisions, in-scope, out-of-scope, success criteria.",
        "Every downstream issue and every line of TDD code inherits the PRD's errors.",
        "Becomes the spec for the whole stack — short PRD, sharp consequences.",
        "Followed immediately by the PRD-REVIEW gate (next slide).",
    ],
    "PRD — the single artifact every downstream stage references.",
)

# 6. ISSUES (to-issues) — vertical slices ────────────────────────────────
stage_slide(
    6, "Stage 3 — Pocock",
    "ISSUES  —  vertical slices.",
    "ISSUES\n(to-issues)",
    AMBER,
    [
        "Breaks the PRD into independently-grabbable vertical slices.",
        "Each slice is end-to-end testable on its own (a real PR, not a layer fragment).",
        "Each slice is sized to fit an LLM's effective context — fresh session, no drift.",
        "`blocked_by:` frontmatter encodes the dependency graph → parallel batches printed automatically.",
        "Slices are reorderable; the PRD stays the contract regardless of which slice ships first.",
    ],
    "ISSUES — the slicing turns one project into many parallelizable units.",
)

# 7. TDD (tdd) ───────────────────────────────────────────────────────────
stage_slide(
    7, "Stage 4 — Pocock",
    "TDD  —  red, green, refactor.",
    "TDD\n(tdd)",
    AMBER,
    [
        "Per slice: write the failing test first, make it pass, then refactor.",
        "Production-code edits are *only* allowed at this stage (PreToolUse hook enforces).",
        "Concurrency-safe commit recipe: explicit pathspecs, never `git add -A`.",
        "Manual-test checklist gate added by Valkyrie — never silently flips `status: done`.",
        "Acceptance criteria checked off as PRs land.",
    ],
    "TDD — the only stage that writes production code.",
)

# 8. AFK (Ralph pattern) ─────────────────────────────────────────────────
stage_slide(
    8, "AFK — Pocock's Ralph pattern",
    "AFK  —  the autonomous loop.",
    "AFK\n(walk away)",
    AMBER,
    [
        "Re-spawns a fresh CLI session per issue — *no context bleed* across slices.",
        "Picks unblocked issues in dependency order automatically.",
        "Budget caps in every dimension: `--max-hours`, `--max-cost-usd`, iteration count.",
        "Closed-stdin abort is a deliberate safety guard against accidental cron launches.",
        "Run from your shell with `!afk N` (you confirm), or `afk N --no-confirm` when Claude launches it.",
    ],
    "AFK — Ralph pattern from aihero.dev, wrapped with budgets and telemetry.",
)

# 9. BLOCK DIAGRAM 1 — Matt Pocock container ─────────────────────────────
s = add_slide()
add_title_bar(s, "The Matt Pocock stack — what we start with.",
              kicker="Block diagram — base")

# Outer Pocock container
add_outline(s, 1.0, 2.0, 11.2, 3.6, "MATT POCOCK — the open-source base",
            color=AMBER, line_w=2.0)

# Four stage blocks inside
stage_w = 2.1
stage_h = 1.5
stage_y = 3.1
gap = 0.25
start_x = 1.4
labels = ["DESIGN\n(grill-with-docs)", "PRD\n(to-prd)",
          "ISSUES\n(to-issues)", "TDD\n(tdd)"]
xs = []
for i, lab in enumerate(labels):
    x = start_x + i * (stage_w + gap + 0.3)
    xs.append(x)
    add_block(s, x, stage_y, stage_w, stage_h, lab, fill=AMBER,
              text=NAVY, size=14)
# arrows between
for i in range(3):
    x1 = xs[i] + stage_w + 0.02
    x2 = xs[i+1] - 0.02
    add_arrow(s, x1, stage_y + stage_h/2, x2, stage_y + stage_h/2,
              color=AMBER)

# AFK wrapper at bottom inside Pocock container
add_outline(s, 1.4, 4.8, 9.8, 0.6, "AFK — Ralph autonomous loop (wraps the slice queue)",
            color=AMBER, line_w=1.5, dash=4)

# Caption
add_text(s, 1.0, 6.0, 11.2, 0.5,
         "DESIGN → PRD → ISSUES → TDD as stages; AFK as the autonomous wrapper.",
         size=14, color=WHITE_DIM, italic=True, align=PP_ALIGN.CENTER)
add_footer(s, "Pocock's pieces, before Valkyrie touches them.", 9, slide_count)

# 10. WHERE POCOCK ALONE BREAKS ──────────────────────────────────────────
s = add_slide()
add_title_bar(s, "Pocock alone, in practice.", kicker="Limits")
add_bullets(s, 0.7, 1.7, 12.0, 5.0, [
    "Walls of text per stage — you re-read three paragraphs to find the next decision.",
    "No mechanical wall — the model can decide it's ready and start writing code.",
    "PRD-REVIEW is honor-based: an enthusiastic \"looks good\" passes.",
    "Long sessions / handoffs lose the orchestrator, gates silently turn into honor-based reminders.",
    "Context limit reached → fresh session restarts without the workflow state.",
    "No place for the *why* (intent) — fuzzy goal in, fuzzy PRD out.",
    "AFK runs use placeholder cost estimates and no real per-iteration audit trail.",
], size=18, line_spacing=1.4)
add_text(s, 0.7, 6.8, 12.0, 0.4,
         "These are the seams Valkyrie was built to close.", size=14,
         color=CYAN, italic=True)
add_footer(s, "Each row above is a Valkyrie feature on the next slides.", 10, slide_count)

# 11. WHAT VALKYRIE ADDS — OVERVIEW ──────────────────────────────────────
s = add_slide()
add_title_bar(s, "What Valkyrie adds.", kicker="Overview")
# Three columns: pre-Design / harness / TDD-side
cols = [
    ("BEFORE DESIGN", VIOLET, [
        "INTENT Lock (new stage)",
        "Optional progressive-enhancement docs",
        "`DOMAIN.md` · `PRODUCT-MAP.md`",
        "`CONTEXT.md` · `docs/intent/*.md`",
        "ADRs",
    ]),
    ("ENFORCEMENT + ERGONOMICS", CYAN, [
        "`/valk` orchestrator skill",
        "Statusline showing live stage",
        "UserPromptSubmit guard (nudge)",
        "PreToolUse TDD gate (hard wall)",
        "PostToolUse telemetry (audit log)",
        "HITL checklists at PRD-REVIEW + TDD completion",
        "Mid-stream loop-back (`valk-revisit`)",
        "Concurrent worktrees (`valk-worktree` / `valk-land`)",
        "Handoff routing rule",
    ]),
    ("AROUND TDD + AFK", EMERALD, [
        "Real cost tracking (OpenAI rates for codex)",
        "Three CLIs: claude, codex, copilot",
        "Escalation ladder: sonnet → opus by default",
        "Manual-test gate before `status: done`",
        "Model-tier rule: insight over code volume",
        "Optional Agent Creator (crew shim)",
    ]),
]
col_w = 4.0
col_gap = 0.2
start_x = 0.55
for i, (head, col, items) in enumerate(cols):
    x = start_x + i * (col_w + col_gap)
    # header band
    add_block(s, x, 1.7, col_w, 0.45, head, fill=col, text=NAVY, size=12)
    # bullets below
    add_bullets(s, x + 0.05, 2.3, col_w - 0.1, 4.6, items, size=12,
                line_spacing=1.35, bullet=" · ")
add_footer(s, "Each item on these lists is a separate, opt-in addition.", 11, slide_count)

# 12. INTENT LOCK (new) ──────────────────────────────────────────────────
s = add_slide()
add_title_bar(s, "INTENT — the gate before the gate.",
              kicker="Valkyrie addition (new stage)")
# Left: a special block, violet, outside Pocock
add_block(s, 0.7, 1.95, 4.0, 4.0, "INTENT LOCK", fill=VIOLET, text=NAVY,
          size=26)
add_bullets(s, 5.2, 2.05, 7.6, 4.0, [
    "Lock the WHY before exploring the HOW.",
    "Every gap is a *question*, never an assumption — \"no inference allowed.\"",
    "Names the domain in one sentence — single repo, bounded context, or umbrella.",
    "Rejects lazy answers (\"build X\" / \"you know what I mean\") with the same pattern as PRD-REVIEW.",
    "Optional persistence to `docs/intent/<slug>.md` (`/to-intent`).",
    "Folds *inside* DESIGN as its opening phase, but conceptually it sits *before* Pocock's pipeline.",
], size=18, line_spacing=1.4)
add_footer(s, "INTENT — colleague's idea: \"intent is the vital first gate.\"", 12, slide_count)

# 13. BLOCK DIAGRAM 2 — INTENT + Pocock + Agent Creator + harness ────────
s = add_slide()
add_title_bar(s, "The full Valkyrie shape.", kicker="Block diagram — extended")

# Pocock container
add_outline(s, 2.4, 2.5, 8.4, 3.3, "MATT POCOCK", color=AMBER, line_w=2.0)
# Stages inside Pocock — smaller blocks
stage_w = 1.6; stage_h = 1.1; stage_y = 3.3
gap = 0.18
inner_start = 2.75
labels = ["DESIGN", "PRD", "ISSUES", "TDD"]
xs = []
for i, lab in enumerate(labels):
    x = inner_start + i * (stage_w + gap + 0.1)
    xs.append(x)
    add_block(s, x, stage_y, stage_w, stage_h, lab, fill=AMBER, text=NAVY,
              size=14)
for i in range(3):
    x1 = xs[i] + stage_w + 0.01
    x2 = xs[i+1] - 0.01
    add_arrow(s, x1, stage_y + stage_h/2, x2, stage_y + stage_h/2,
              color=AMBER)
# AFK wrapper inside Pocock
add_outline(s, 2.75, 4.7, 7.7, 0.5, "AFK — autonomous loop", color=AMBER,
            line_w=1.2, dash=4)

# INTENT block — left, OUTSIDE Pocock
add_block(s, 0.55, 3.3, 1.4, 1.1, "INTENT", fill=VIOLET, text=NAVY, size=14)
add_arrow(s, 1.95, 3.85, 2.75, 3.85, color=VIOLET)
add_text(s, 0.55, 4.5, 1.4, 0.3, "Valkyrie", size=10, color=VIOLET,
         italic=True, align=PP_ALIGN.CENTER)

# Agent Creator block — right, OUTSIDE Pocock, attaches to TDD
add_block(s, 11.0, 3.3, 1.7, 1.1, "AGENT\nCREATOR\n(optional)", fill=EMERALD,
          text=NAVY, size=12)
add_arrow(s, xs[3] + stage_w + 0.02, 3.85, 11.0 - 0.02, 3.85, color=EMERALD)
add_text(s, 11.0, 4.5, 1.7, 0.3, "Valkyrie · crew shim", size=10,
         color=EMERALD, italic=True, align=PP_ALIGN.CENTER)

# Harness band at bottom — full width
add_outline(s, 0.55, 5.95, 12.2, 1.0, "VALKYRIE HARNESS  —  /valk · statusline · 3 hooks · valk-revisit · valk-worktree · valk-land  ·  HITL gates",
            color=CYAN, line_w=1.5)

add_footer(s, "Pocock at the core; Valkyrie wraps and extends.", 13, slide_count)

# 14. THE HARNESS — statusline + 3 hooks ─────────────────────────────────
s = add_slide()
add_title_bar(s, "The harness — what makes it stick.", kicker="Enforcement layer")
# 4 columns, one per harness component
items = [
    ("STATUSLINE\n▶ STAGE", CYAN,
     "Every prompt shows the current stage.\n\nYou always know whether the next message is design talk or code.",
     "statusline.py"),
    ("USER-PROMPT\nGUARD", CYAN,
     "Nudges build prompts into `/valk` instead of jumping to code.\n\nSoft — model can still ignore.",
     "valk-guard.sh"),
    ("TDD GATE\n(hard wall)", EMERALD,
     "PreToolUse hook.\n\nLiterally DENIES `Edit` / `Write` / `Bash` writes to src/ unless stage is `tdd`.\n\nNot honor-based — mechanical.",
     "valk-tdd-gate.sh"),
    ("TELEMETRY\n(audit log)", CYAN,
     "PostToolUse hook.\n\nLogs every Read/Edit during active stages — the AFK audit trail.",
     "valk-telemetry.sh"),
]
col_w = 2.95
col_gap = 0.18
start_x = 0.55
for i, (head, col, body, fname) in enumerate(items):
    x = start_x + i * (col_w + col_gap)
    add_block(s, x, 1.8, col_w, 1.1, head, fill=col, text=NAVY, size=14)
    add_text(s, x + 0.1, 3.05, col_w - 0.2, 3.5, body, size=12,
             color=WHITE_DIM)
    add_text(s, x + 0.1, 6.45, col_w - 0.2, 0.4, fname, size=10,
             color=CYAN_DIM, italic=True, font="Menlo")
add_footer(s, "Three hooks + a statusline = enforcement that survives long sessions.",
           14, slide_count)

# 15. THE TDD GATE — mechanical wall visualization ───────────────────────
s = add_slide()
add_title_bar(s, "The mechanical wall.", kicker="The load-bearing hook")
# Left side: stages with red X for "no code"
add_text(s, 0.7, 1.85, 5, 0.4, "STAGE", size=12, color=CYAN_DIM, bold=True)
stages = [
    ("idle",       True,  "you're not in a flow"),
    ("design",     False, "no production-code edits"),
    ("prd",        False, "no production-code edits"),
    ("prd-review", False, "no production-code edits"),
    ("issues",     False, "no production-code edits"),
    ("tdd",        True,  "production code allowed"),
    ("afk",        True,  "loop runs TDD per slice"),
]
ay = 2.3
for st, allow, note in stages:
    color = EMERALD if allow else ROSE
    mark = "✓" if allow else "✕"
    # stage name
    add_text(s, 0.7, ay, 2.0, 0.4, st, size=15, color=WHITE)
    # marker
    add_text(s, 2.8, ay, 0.5, 0.4, mark, size=18, color=color, bold=True)
    # note
    add_text(s, 3.3, ay + 0.04, 4.5, 0.4, note, size=12, color=WHITE_DIM, italic=True)
    ay += 0.45

# Right side: how it works diagram
add_text(s, 8.4, 1.85, 4.5, 0.4, "HOW IT WORKS", size=12, color=CYAN_DIM, bold=True)
add_block(s, 8.4, 2.35, 4.4, 0.7, "Model tries to Edit src/foo.ts", fill=PANEL, text=WHITE,
          size=12, bold=False, line=BORDER)
add_arrow(s, 10.6, 3.07, 10.6, 3.45, color=WHITE_DIM)
add_block(s, 8.4, 3.5, 4.4, 0.7, "PreToolUse hook fires (valk-tdd-gate.sh)",
          fill=PANEL, text=WHITE, size=12, bold=False, line=BORDER)
add_arrow(s, 10.6, 4.22, 10.6, 4.6, color=WHITE_DIM)
add_block(s, 8.4, 4.65, 4.4, 0.7, "stage.py get  →  reads .claude/valk/stage",
          fill=PANEL, text=WHITE, size=12, bold=False, line=BORDER)
add_arrow(s, 10.6, 5.37, 10.6, 5.75, color=WHITE_DIM)
add_block(s, 8.4, 5.8, 4.4, 0.7, "if stage != tdd → DENY", fill=ROSE, text=WHITE,
          size=14, bold=True)
add_text(s, 8.4, 6.65, 4.4, 0.3, "threat model: model drift, not sabotage",
         size=10, color=WHITE_FAINT, italic=True, align=PP_ALIGN.CENTER)
add_footer(s, "Mechanical, not honor-based. Markdown stays writable everywhere.",
           15, slide_count)

# 16. CONCISE UI — bullets + options vs prose ────────────────────────────
s = add_slide()
add_title_bar(s, "Concise UI > walls of text.",
              kicker="Prompt-injection improvements")
# left: raw Pocock-style
add_text(s, 0.7, 1.65, 6.0, 0.4, "RAW: prose monologue", size=12,
         color=ROSE, bold=True)
prose = ("\"Now that we've finished gathering requirements, we can move on "
         "to the PRD. The PRD should contain the in-scope items, the out-of-"
         "scope items, the implementation decisions, the open questions, "
         "the success criteria, and the rationale for each implementation "
         "decision. Once you're satisfied with the PRD, we can move on to "
         "issue breakdown. Are you ready?\"")
tb = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.7), Inches(2.1), Inches(6.0), Inches(4.7))
tb.fill.solid(); tb.fill.fore_color.rgb = PANEL; tb.line.color.rgb = BORDER
tb.line.width = Pt(1)
try: tb.adjustments[0] = 0.05
except Exception: pass
tf = tb.text_frame
tf.margin_left = tf.margin_right = Inches(0.2)
tf.margin_top = tf.margin_bottom = Inches(0.2)
tf.word_wrap = True
p = tf.paragraphs[0]
r = p.add_run()
r.text = prose
r.font.size = Pt(14); r.font.color.rgb = WHITE_DIM; r.font.italic = True
r.font.name = "Helvetica Neue"
# right: Valkyrie-style options
add_text(s, 7.3, 1.65, 5.5, 0.4, "VALKYRIE: structured options",
         size=12, color=EMERALD, bold=True)
# Question
add_text(s, 7.3, 2.1, 5.5, 1.0,
         "All 5 decisions above are in scope and will all be built.",
         size=13, color=WHITE, bold=True)
add_block(s, 7.3, 3.3, 5.5, 0.65,
          "Approve — build all 5 decisions (full scope)",
          fill=EMERALD, text=NAVY, size=13)
add_block(s, 7.3, 4.1, 5.5, 0.65, "Redline a decision",
          fill=AMBER, text=NAVY, size=13)
add_block(s, 7.3, 4.9, 5.5, 0.65, "Open the file first", fill=PANEL,
          text=WHITE, size=13, line=BORDER)
add_text(s, 7.3, 5.85, 5.5, 1.0,
         "Same decision; bare \"yes\" / \"lgtm\" still rejected — must click one.",
         size=11, color=WHITE_FAINT, italic=True)
add_footer(s, "Every gate works this way: clear options, clear refusal pattern.",
           16, slide_count)

# 17. HITL — blockers and enforcement ────────────────────────────────────
s = add_slide()
add_title_bar(s, "Human-in-the-loop gates.", kicker="Blockers")
gates = [
    ("PRD-REVIEW", VIOLET,
     "Single \"Approve all N decisions\" option.\n\"yes\" / \"lgtm\" rejected.\nHuman must click — every time, this conversation.",
     "highest-leverage gate"),
    ("TDD MANUAL\nCHECKLIST", CYAN,
     "Before flipping `status: done`, parse the issue for `## Manual test checklist`.\nUnticked items → halt + pass/fail/wait prompt.",
     "real trial caught a UX redesign here"),
    ("MID-STREAM\nLOOP-BACK", EMERALD,
     "`/valk` watches for change signals at every downstream stage.\nDetected → single confirm → `valk-revisit` runs.\nArtifacts amended in place.",
     "TDD never `rm`s tests"),
    ("HANDOFF\nROUTING", AMBER,
     "If the user `/handoff`s mid-flow, the doc MUST instruct: \"Run /valk first.\"\nNever script direct skill invocations.",
     "stops silent gate-skipping"),
]
col_w = 2.95
col_gap = 0.18
start_x = 0.55
for i, (head, col, body, foot) in enumerate(gates):
    x = start_x + i * (col_w + col_gap)
    add_block(s, x, 1.8, col_w, 1.1, head, fill=col, text=NAVY, size=13)
    add_text(s, x + 0.1, 3.05, col_w - 0.2, 3.0, body, size=12,
             color=WHITE_DIM)
    add_text(s, x + 0.1, 6.5, col_w - 0.2, 0.35, foot, size=10,
             color=col, italic=True, bold=True)
add_footer(s, "Four explicit gates — each refusable, each humanly enforceable.",
           17, slide_count)

# 18. AFK ENHANCEMENTS ───────────────────────────────────────────────────
s = add_slide()
add_title_bar(s, "AFK — what Valkyrie added on top of Ralph.",
              kicker="Autonomous loop, hardened")
add_bullets(s, 0.7, 1.7, 6.0, 5.5, [
    "Three `--cli` backends: claude, codex, copilot.",
    "Each CLI uses its own login — afk never handles API keys.",
    "Real cost tracking for codex via `codex exec --json` + OpenAI rates.",
    "Budget caps in every dimension: --max-hours, --max-cost-usd, iter count.",
    "Closed-stdin abort = deliberate safety guard against accidental launches.",
], size=15, line_spacing=1.4)
add_bullets(s, 7.0, 1.7, 6.0, 5.5, [
    "Escalation ladder: sonnet → opus on failure (default-on, claude only).",
    "Non-claude cost-parse miss is non-fatal — never kills the loop.",
    "Per-iteration cost CSV: `.claude/valk/afk-cost-history.csv`.",
    "Per-iteration log files: `.claude/valk/afk-logs/`.",
    "Confirmation queue printed up-front: caps, queue order, dependency reasoning.",
    "`!afk N` (user-driven) vs `afk N --no-confirm` (Claude-launched).",
], size=15, line_spacing=1.4)
add_footer(s, "Same Ralph idea; now with budgets, telemetry, and three model providers.",
           18, slide_count)

# 19. MODEL TIER — INSIGHT, NOT CODE VOLUME ──────────────────────────────
s = add_slide()
add_title_bar(s, "Model tier follows leverage, not code volume.",
              kicker="Cost discipline")
# Two-row table: stage, recommended tier, reason
rows = [
    ("DESIGN / PRD / ISSUES", "strongest tier",
     "broad context · judgment · refusal to settle for fuzzy answers",
     VIOLET),
    ("TDD",                   "sonnet",
     "pattern-matching against a clear spec; escalation ladder starts here",
     CYAN),
    ("Delegated reads / QA",  "haiku",
     "single-task agents on short investigations",
     EMERALD),
]
ay = 2.2
add_text(s, 0.7,  1.9, 3.4, 0.4, "STAGE",     size=11, color=CYAN_DIM, bold=True)
add_text(s, 4.6,  1.9, 2.4, 0.4, "TIER",      size=11, color=CYAN_DIM, bold=True)
add_text(s, 7.4,  1.9, 5.4, 0.4, "WHY",       size=11, color=CYAN_DIM, bold=True)
for stg, tier, why, col in rows:
    add_block(s, 0.7,  ay, 3.6, 0.85, stg, fill=col, text=NAVY, size=14)
    add_block(s, 4.5,  ay, 2.6, 0.85, tier, fill=PANEL, text=WHITE, size=14, line=BORDER)
    add_text (s, 7.3,  ay + 0.18, 5.5, 0.6, why, size=12, color=WHITE_DIM)
    ay += 1.05

# Anti-pattern callout
add_outline(s, 0.7, 5.7, 12.0, 1.05, None, color=ROSE, line_w=1.5)
add_text(s, 0.9, 5.85, 1.0, 0.4, "⚠", size=22, color=ROSE)
add_text(s, 1.5, 5.78, 11.0, 0.5,
         "Anti-pattern: \"no code in this stage = cheap model.\"",
         size=14, color=ROSE, bold=True)
add_text(s, 1.5, 6.18, 11.0, 0.5,
         "Inverts the leverage. Cost of a bad plan dwarfs the cost of a strong model "
         "on the conversation that produces it.", size=12, color=WHITE_DIM)
add_footer(s, "Was internally contradictory in the SKILL.md before this fix.",
           19, slide_count)

# 20. PROGRESSIVE-ENHANCEMENT DOCS — the optional .md set ────────────────
s = add_slide()
add_title_bar(s, "The optional docs.", kicker="Progressive enhancement")
docs_list = [
    ("INTENT", VIOLET, "docs/intent/<slug>.md", "per task",
     "the why, in/out-of-scope, success criteria"),
    ("DOMAIN", AMBER, "DOMAIN.md", "per repo",
     "bounds, integrations, installer relationship, legacy constraints"),
    ("PRODUCT-MAP", AMBER, "PRODUCT-MAP.md", "umbrella",
     "member repos, build order, cross-repo contracts"),
    ("CONTEXT", CYAN, "CONTEXT.md (or CONTEXT-MAP.md + per-context)",
     "per repo / per bounded context", "glossary — terms + relationships"),
    ("ADRs", EMERALD, "docs/adr/*.md", "per decision",
     "architectural memory: hard-to-reverse + non-obvious + real trade-off"),
]
ay = 1.85
add_text(s, 0.7, ay, 1.6, 0.35, "DOC",   size=11, color=CYAN_DIM, bold=True)
add_text(s, 2.4, ay, 4.0, 0.35, "FILE",  size=11, color=CYAN_DIM, bold=True)
add_text(s, 6.5, ay, 1.8, 0.35, "SCOPE", size=11, color=CYAN_DIM, bold=True)
add_text(s, 8.5, ay, 4.3, 0.35, "HOLDS", size=11, color=CYAN_DIM, bold=True)
ay += 0.4
for name, col, path, scope, holds in docs_list:
    add_block(s, 0.7, ay, 1.6, 0.6, name, fill=col, text=NAVY, size=12)
    add_text (s, 2.4, ay + 0.16, 4.0, 0.4, path, size=11, color=WHITE, font="Menlo")
    add_text (s, 6.5, ay + 0.16, 1.8, 0.4, scope, size=11, color=WHITE_DIM, italic=True)
    add_text (s, 8.5, ay + 0.16, 4.3, 0.4, holds, size=11, color=WHITE_DIM)
    ay += 0.75

add_text(s, 0.7, 6.65, 12.0, 0.4,
         "All optional. Every skill is a no-op against absent files.  \"Enforcement scales with what you write.\"",
         size=13, color=CYAN, italic=True, align=PP_ALIGN.CENTER)
add_footer(s, "Each doc adds one upstream gate; absent files = today's behaviour.",
           20, slide_count)

# 21. WHERE EACH DOC IS USED (matrix) ────────────────────────────────────
s = add_slide()
add_title_bar(s, "Where each doc is read.", kicker="Progressive enhancement · matrix")
# stages across, docs down
stages = ["INTENT", "DESIGN", "PRD", "ISSUES", "TDD"]
docs = [
    ("intent.md",      [False, True,  True,  False, False],
                       "PRD must trace back to the locked why"),
    ("DOMAIN.md",      [True,  True,  True,  False, False],
                       "grilling grounded in bounds; PRD cites bounds"),
    ("PRODUCT-MAP.md", [True,  True,  True,  False, False],
                       "cross-repo contract checks"),
    ("CONTEXT.md",     [False, True,  False, False, False],
                       "glossary challenges; written inline as terms resolve"),
    ("docs/adr/*.md",  [False, True,  True,  True,  False],
                       "new ADRs written; downstream stages cite them"),
]
# header row
col_x = [3.5, 4.9, 6.2, 7.5, 8.8]
col_w_each = 1.2
header_y = 1.85
add_text(s, 0.6, header_y, 2.8, 0.4, "DOC", size=12, color=CYAN_DIM, bold=True)
for i, st in enumerate(stages):
    add_text(s, col_x[i], header_y, col_w_each, 0.4, st, size=11, color=CYAN_DIM, bold=True, align=PP_ALIGN.CENTER)
add_text(s, 10.3, header_y, 2.8, 0.4, "ROLE", size=12, color=CYAN_DIM, bold=True)

row_y = 2.4
for name, hits, role in docs:
    add_text(s, 0.6, row_y, 2.8, 0.45, name, size=12, color=WHITE, font="Menlo")
    for i, hit in enumerate(hits):
        if hit:
            add_block(s, col_x[i], row_y + 0.03, 0.7, 0.4, "●", fill=CYAN, text=NAVY, size=12, bold=True)
        else:
            add_text(s, col_x[i], row_y + 0.05, 1.2, 0.4, "—", size=12, color=WHITE_FAINT, align=PP_ALIGN.CENTER)
    add_text(s, 10.3, row_y + 0.05, 2.8, 0.6, role, size=10, color=WHITE_DIM, italic=True)
    row_y += 0.75

add_text(s, 0.6, 6.6, 12.0, 0.4,
         "Cumulative: each new file you author strengthens every upstream gate.",
         size=13, color=CYAN, italic=True)
add_footer(s, "ISSUES/TDD don't read these directly — but inherit their effects via PRD.",
           21, slide_count)

# 22. LOOP-BACK + WORKTREES — flow-shape tooling ─────────────────────────
s = add_slide()
add_title_bar(s, "Real-engineering ergonomics.",
              kicker="Mid-stream change + concurrent flows")
# Two side-by-side cards
def feature_card(x, y, w, h, head, sub, items, color):
    add_block(s, x, y, w, 0.7, head, fill=color, text=NAVY, size=14)
    add_text(s, x + 0.05, y + 0.78, w - 0.1, 0.35, sub,
             size=11, color=color, italic=True, bold=True)
    add_bullets(s, x + 0.05, y + 1.2, w - 0.1, h - 1.3, items, size=12,
                color=WHITE_DIM, line_spacing=1.35, bullet=" · ")

feature_card(0.55, 1.85, 6.1, 5.0, "LOOP-BACK", "valk-revisit",
             [
                 "/valk watches for change signals every downstream stage.",
                 "Single confirmation prompt — never silent.",
                 "valk-revisit writes docs/changes/<ts>-<slug>.md.",
                 "Rewinds the stage marker; re-entered skill amends artifacts in place.",
                 "Δ change-log entries preserve the trail.",
                 "Hard rule: TDD loop-back never `rm`s tests.",
             ], EMERALD)

feature_card(6.85, 1.85, 6.1, 5.0, "CONCURRENT FLOWS",
             "valk-worktree · valk-land",
             [
                 "One Valkyrie flow per terminal — no shared stage marker.",
                 "valk-worktree <name>: linked worktree on valk/<name> branch.",
                 "Each flow runs its own DESIGN → PRD → ISSUES → TDD.",
                 "Independent issue batches run in parallel (one worktree each).",
                 "valk-land <name>: race-free integrate-back to main.",
                 "Lets one human (or one AFK loop) ship multiple slices at once.",
             ], CYAN)

add_footer(s, "Requirements change · multiple slices in flight · no \"resume from scratch\".",
           22, slide_count)

# 23. CREW SHIM — Agent Creator ──────────────────────────────────────────
s = add_slide()
add_title_bar(s, "Crew Shim — Agent Creator at TDD (optional).",
              kicker="Drop-in agent composition")
add_text(s, 0.7, 1.7, 12.0, 0.45,
         ".claude/valk-config.md binds a forged crew to one or more stages.",
         size=15, color=EMERALD, italic=True)
add_text(s, 0.7, 2.2, 12.0, 0.45,
         "Each stage's vanilla skill becomes \"dispatch to bound crew\" instead.",
         size=15, color=EMERALD, italic=True)

# Diagram: stage skill → vanilla OR crew dispatch
add_block(s, 0.7,  3.2, 2.6, 1.0, "Stage skill\n(e.g. /tdd)",
          fill=AMBER, text=NAVY, size=13)
add_arrow(s, 3.4, 3.7, 4.1, 3.7, color=WHITE_DIM)
add_block(s, 4.2,  3.2, 2.4, 1.0, "crew-shim decide",
          fill=PANEL, text=WHITE, size=13, line=BORDER)
add_arrow(s, 6.7, 3.45, 7.5, 3.0, color=WHITE_DIM)
add_arrow(s, 6.7, 3.95, 7.5, 4.4, color=WHITE_DIM)
add_block(s, 7.55, 2.5, 5.0, 1.0,
          "VANILLA — run stage's stock skill",
          fill=PANEL, text=WHITE, size=12, line=BORDER)
add_block(s, 7.55, 3.95, 5.0, 1.5,
          "CREW — dispatch bound agents per task\n(implementer → tester → reviewer-gate)",
          fill=EMERALD, text=NAVY, size=12)

add_bullets(s, 0.7, 5.4, 12.0, 1.5, [
    "Stage enforcement (order, PRD-REVIEW gate, TDD wall) is unchanged — only the *worker* swaps.",
    "Default = vanilla. No config = no shim. Byte-identical to having no shim at all.",
    "A gating agent (reviewer/security) writing `status: blocked` halts the line.",
], size=13, line_spacing=1.35)
add_footer(s, "Authored crews from agent-builder forge plug in here.",
           23, slide_count)

# 24. THE COMPOUNDING EFFECT ─────────────────────────────────────────────
s = add_slide()
add_title_bar(s, "Compounding effect.", kicker="Why each piece pays for itself")
# Stack of bars showing each addition adding strength
items = [
    ("Just Pocock",                     "stages exist, but every gate is honor-based",      1.5, AMBER),
    ("+ CONTEXT.md",                    "terms challenged at DESIGN, fuzzy language sharpened", 2.5, CYAN),
    ("+ DOMAIN.md",                     "grilling refuses to wander out of bounds",          3.5, CYAN),
    ("+ PRODUCT-MAP.md",                "cross-repo contracts named before any code",        4.5, CYAN),
    ("+ intent.md",                     "PRD must trace back to a locked why",               5.5, CYAN),
    ("+ ADRs",                          "expensive decisions outlive the conversation",      6.5, EMERALD),
    ("+ Valkyrie harness (hooks, gates, /valk)", "mechanical wall + HITL gates + observability", 8.0, EMERALD),
]
for label, why, length, col in items:
    add_block(s, 0.6, 1.85 + items.index((label,why,length,col)) * 0.65,
              length, 0.5, label, fill=col, text=NAVY, size=12)
    add_text(s, length + 0.8, 1.95 + items.index((label,why,length,col)) * 0.65,
             8, 0.4, why, size=12, color=WHITE_DIM, italic=True)
add_text(s, 0.7, 6.7, 12.0, 0.4,
         "Each row above is opt-in. Don't write the docs you don't want — but each one strengthens every gate above it.",
         size=11, color=WHITE_FAINT, italic=True, align=PP_ALIGN.CENTER)
add_footer(s, "\"Enforcement scales with what you write.\"", 24, slide_count)

# 25. TAKEAWAY ───────────────────────────────────────────────────────────
s = add_slide()
add_title_bar(s, "Takeaway.", kicker="The bet")
add_text(s, 0.7, 1.9, 12.0, 0.7, "Pocock gave us the skill primitives.",
         size=26, color=AMBER, bold=True)
add_text(s, 0.7, 2.8, 12.0, 0.7,
         "Valkyrie adds enforcement, gates, observability, ergonomics, telemetry.",
         size=24, color=CYAN, bold=True)
add_text(s, 0.7, 4.0, 12.0, 2.5,
         "The workflow is the product — not the model.\n\n"
         "We adopt the best community skills, then layer mechanical guarantees on top.\n"
         "Every Valkyrie addition is opt-in. Every addition compounds with the others.\n"
         "Built incrementally from real dogfooding — feature added every time the workflow leaked.",
         size=16, color=WHITE_DIM)
add_footer(s, "What you actually take to your team.", 25, slide_count)

# 26. ROADMAP ────────────────────────────────────────────────────────────
s = add_slide()
add_title_bar(s, "Roadmap.", kicker="What's next")
items = [
    ("Auto-PR generation via Azure DevOps", VIOLET,
     "`pr_skill: to-azure-pr` already plugs into /tdd's completion gate.\nNext: write the skill itself.\nLinked work-item, branch convention, CI wait."),
    ("Auto-documentation on commit", EMERALD,
     "PostCommit skill that updates DOMAIN.md / CONTEXT.md / READMEs whenever code lands.\nCodebase docs stay current automatically."),
    ("Org-level DOMAIN / CONTEXT / PRODUCT-MAP lockdown", CYAN,
     "Org-wide templates for the three load-bearing docs.\nNew repos get the skeleton on first /valk."),
    ("Profile system (team-specific variants)", AMBER,
     "profiles/azure-devops, profiles/<team>/ layered on the baseline.\nKeeps Valkyrie tool-agnostic by default."),
]
ay = 1.85
for head, col, body in items:
    add_block(s, 0.7, ay, 4.0, 1.05, head, fill=col, text=NAVY, size=14)
    add_text(s, 4.9, ay + 0.1, 8.0, 1.0, body, size=12, color=WHITE_DIM,
             )
    ay += 1.25
add_footer(s, "Each item is a layer on the existing base — same opt-in spirit.",
           26, slide_count)

# 27. OTHER PIECES YOU MIGHT'VE MISSED ───────────────────────────────────
s = add_slide()
add_title_bar(s, "Other pieces worth knowing.", kicker="Survey of the rest")
add_bullets(s, 0.55, 1.7, 6.2, 5.5, [
    "valk-worktree / valk-land — concurrent flows on one checkout.",
    "valk-revisit — mid-stream loop-back (the script behind /valk's auto-detection).",
    "valk-tdd-gate.sh — the PreToolUse mechanical wall.",
    "valk-guard.sh — UserPromptSubmit nudge into /valk.",
    "valk-telemetry.sh — PostToolUse audit log.",
    "/zoom-out — escape-hatch skill to re-orient on unfamiliar code.",
    "/refactor-spaghetti — find deepening opportunities in tangled code.",
], size=13, color=WHITE_DIM, line_spacing=1.4)
add_bullets(s, 6.95, 1.7, 6.2, 5.5, [
    "install.sh — idempotent installer with symlink-or-copy fallback (Git Bash).",
    "Per-skill SKILL.md verification echo (\"valk: synced 12,893 bytes\").",
    "afk telemetry: per-iteration log + cost CSV under .claude/valk/.",
    "Auto cost-mode selection (subscription vs API-billed).",
    "Stage helper (stage.py) + statusline (statusline.py).",
    "test/run-tests.sh: 16 tests; every Valkyrie addition has coverage.",
    "Handoff routing rule (mid-Valkyrie /handoff must say \"Run /valk first\").",
], size=13, color=WHITE_DIM, line_spacing=1.4)
add_footer(s, "All of this in github.com/BunnyDAO/Valkyrie — docs/ has the SOP + workflow walkthroughs.",
           27, slide_count)

# 28. Q & A ───────────────────────────────────────────────────────────────
s = add_slide()
# big centered Q&A
add_text(s, 0.6, 2.8, 12.0, 1.2, "Q & A", size=72, bold=True, color=CYAN,
         align=PP_ALIGN.CENTER)
add_text(s, 0.6, 4.2, 12.0, 0.7,
         "github.com/BunnyDAO/Valkyrie",
         size=22, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, 0.6, 4.95, 12.0, 0.5,
         "skills · scripts · hooks · docs · SOP", size=14,
         color=WHITE_DIM, italic=True, align=PP_ALIGN.CENTER)
# subtle rule
r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.0), Inches(5.8),
                      Inches(1.3), Inches(0.04))
r.fill.solid(); r.fill.fore_color.rgb = CYAN; r.line.fill.background()

# ─── save ─────────────────────────────────────────────────────────────────
out = Path(__file__).resolve().parent / "valkyrie-overview.pptx"
prs.save(out)
print(f"Wrote {out}")
print(f"Slides: {len(prs.slides)}")
